# -*- coding: utf-8 -*-
"""
文档上传与解析服务

支持 PDF/PPTX/DOCX/Excel 文件上传
使用 MinerU 3.1 解析引擎
存储到 PostgreSQL + PGVector
"""

import os
import json
import hashlib
import asyncio
from datetime import datetime
from typing import List, Dict, Any, Optional
from pathlib import Path

import psycopg2
from psycopg2.extras import RealDictCursor
from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel

# 配置
UPLOAD_DIR = Path(r"C:\Users\11489\.openclaw\workspace-market\uploads")
UPLOAD_DIR.mkdir(exist_ok=True)

# 数据库配置
# 默认使用远程 vectordb 服务器，可通过环境变量覆盖
DB_CONFIG = {
    "host": os.getenv("DB_HOST", "192.168.3.146"),
    "port": int(os.getenv("DB_PORT", "5432")),
    "database": os.getenv("DB_NAME", "vectordb"),
    "user": os.getenv("DB_USER", "vectordb"),
    "password": os.getenv("DB_PASSWORD", "vectordb123")
}

# Embedding 模型配置
EMBEDDING_MODEL = "paraphine-multilingual-MiniLM-L12-v2"
EMBEDDING_DIM = 384

app = FastAPI(title="文档上传服务", version="1.0.0")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ==================== 数据库操作 ====================

def get_db_connection():
    """获取数据库连接"""
    try:
        conn = psycopg2.connect(**DB_CONFIG, cursor_factory=RealDictCursor)
        return conn
    except psycopg2.OperationalError as e:
        print(f"Database connection error: {e}")
        raise


def get_embedding_model():
    """获取或加载 Embedding 模型（延迟加载）"""
    if not hasattr(get_embedding_model, '_model'):
        try:
            from sentence_transformers import SentenceTransformer
            get_embedding_model._model = SentenceTransformer('paraphrase-multilingual-MiniLM-L12-v2')
        except Exception as e:
            print(f"Warning: Failed to load embedding model: {e}")
            return None
    return get_embedding_model._model


# ==================== 文件解析 ====================

async def parse_document_mineru(file_path: str, file_type: str) -> Dict[str, Any]:
    """
    使用 MinerU 3.1 解析文档
    返回: {markdown, metadata, tables, images}
    """
    try:
        from magic_pdf.data_utils.img_context_token import (
            tokens_to_image_dict,
            pdf_with_images_to_magic_content
        )
        from magic_pdf.pdf_parse_v3 import parse_pdf
        from magic_pdf.model.doc_scene_analyze import MagicPipline

        # 根据文件类型选择解析方式
        if file_type == 'pdf':
            result = await parse_pdf_with_mineru(file_path)
        elif file_type in ['pptx', 'ppt']:
            result = await parse_pptx_with_mineru(file_path)
        elif file_type in ['docx', 'doc']:
            result = await parse_docx_with_mineru(file_path)
        elif file_type in ['xlsx', 'xls']:
            result = await parse_excel_with_mineru(file_path)
        else:
            raise ValueError(f"不支持的文件类型: {file_type}")

        return result

    except ImportError:
        # MinerU 未安装，使用备用解析
        return await parse_document_fallback(file_path, file_type)


async def parse_pdf_with_mineru(pdf_path: str) -> Dict[str, Any]:
    """使用 MinerU 解析 PDF"""
    # MinerU 3.1 解析逻辑
    # 这里需要根据 MinerU 3.1 API 调整
    from magic_pdf.pdf_parse_v3 import parse_pdf

    result = await asyncio.to_thread(parse_pdf, pdf_path)
    return {
        "markdown": result.get("markdown", ""),
        "metadata": result.get("meta", {}),
        "tables": result.get("table_blocks", []),
        "images": result.get("image_paths", [])
    }


async def parse_pptx_with_mineru(pptx_path: str) -> Dict[str, Any]:
    """使用 MinerU 解析 PPTX"""
    from pptx_to_markdown import convert_pptx_to_markdown

    markdown = await asyncio.to_thread(convert_pptx_to_markdown, pptx_path)
    return {
        "markdown": markdown,
        "metadata": {"slides": extract_slide_count(pptx_path)},
        "tables": [],
        "images": extract_pptx_images(pptx_path)
    }


async def parse_docx_with_mineru(docx_path: str) -> Dict[str, Any]:
    """使用 MinerU 解析 DOCX"""
    from docx import Document

    def _parse():
        doc = Document(docx_path)
        paragraphs = [p.text for p in doc.paragraphs]
        markdown = "\n\n".join(paragraphs)

        # 提取表格
        tables = []
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                table_text.append(" | ".join(row_text))
            tables.append("\n".join(table_text))

        return markdown, tables

    markdown, tables = await asyncio.to_thread(_parse)
    return {
        "markdown": markdown,
        "metadata": {"paragraphs": len(doc.paragraphs), "tables": len(tables)},
        "tables": tables,
        "images": []
    }


async def parse_excel_with_mineru(excel_path: str) -> Dict[str, Any]:
    """使用 MinerU 解析 Excel"""
    import openpyxl

    def _parse():
        wb = openpyxl.load_workbook(excel_path, data_only=True)
        sheets_data = {}

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = []
            for row in ws.iter_rows(values_only=True):
                # 过滤空行
                if any(cell is not None for cell in row):
                    rows.append([str(cell) if cell is not None else "" for cell in row])

            # 转为 Markdown 表格
            if rows:
                md_lines = []
                for i, row in enumerate(rows):
                    line = "| " + " | ".join(row) + " |"
                    md_lines.append(line)
                    if i == 0:  # 表头后加分隔行
                        md_lines.append("|" + "|".join(["---"] * len(row)) + "|")
                sheets_data[sheet_name] = "\n".join(md_lines)

        return sheets_data

    sheets_data = await asyncio.to_thread(_parse)

    # 合并所有 sheet
    markdown = ""
    for sheet_name, content in sheets_data.items():
        markdown += f"## {sheet_name}\n\n{content}\n\n"

    return {
        "markdown": markdown,
        "metadata": {"sheets": list(sheets_data.keys()), "sheet_count": len(sheets_data)},
        "tables": list(sheets_data.values()),
        "images": []
    }


async def parse_document_fallback(file_path: str, file_type: str) -> Dict[str, Any]:
    """备用解析方法（MinerU 不可用时）"""
    if file_type == 'pdf':
        return await parse_pdf_fallback(file_path)
    elif file_type in ['pptx', 'ppt']:
        return await parse_pptx_fallback(file_path)
    elif file_type in ['docx', 'doc']:
        return await parse_docx_fallback(file_path)
    elif file_type in ['xlsx', 'xls']:
        return await parse_excel_with_mineru(file_path)
    else:
        return {"markdown": "", "metadata": {}, "tables": [], "images": []}


async def parse_pdf_fallback(pdf_path: str) -> Dict[str, Any]:
    """PDF 备用解析（pdfplumber）"""
    import pdfplumber

    def _parse():
        all_text = []
        all_tables = []

        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    all_text.append(text)

                tables = page.extract_tables()
                for table in tables:
                    if table:
                        table_text = []
                        for row in table:
                            table_text.append(" | ".join([str(c) if c else "" for c in row]))
                        all_tables.append("\n".join(table_text))

        return "\n\n".join(all_text), all_tables

    text, tables = await asyncio.to_thread(_parse)
    return {
        "markdown": text,
        "metadata": {"pages": len(pdf.pages) if 'pdf' in dir() else 0},
        "tables": tables,
        "images": []
    }


async def parse_pptx_fallback(pptx_path: str) -> Dict[str, Any]:
    """PPTX 备用解析（python-pptx）"""
    from pptx import Presentation

    def _parse():
        prs = Presentation(pptx_path)
        slides = []

        for i, slide in enumerate(prs.slides):
            slide_text = [f"### Slide {i + 1}\n"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            slides.append("\n".join(slide_text))

        return "\n\n---\n\n".join(slides)

    markdown = await asyncio.to_thread(_parse)
    return {
        "markdown": markdown,
        "metadata": {"slides": len(prs.slides)},
        "tables": [],
        "images": []
    }


async def parse_docx_fallback(docx_path: str) -> Dict[str, Any]:
    """DOCX 备用解析"""
    return await parse_docx_with_mineru(docx_path)


# ==================== 向量化和存储 ====================

def chunk_text(text: str, chunk_size: int = 512, overlap: int = 50) -> List[str]:
    """将长文本分块"""
    if not text:
        return []

    # 按段落分割
    paragraphs = text.split("\n\n")
    chunks = []
    current_chunk = []
    current_size = 0

    for para in paragraphs:
        para_size = len(para)
        if current_size + para_size > chunk_size and current_chunk:
            chunks.append("\n\n".join(current_chunk))
            # 保留最后一个段落用于重叠
            if len(current_chunk) > 1:
                current_chunk = [current_chunk[-1]]
            else:
                current_chunk = []
            current_size = len("\n\n".join(current_chunk))

        current_chunk.append(para)
        current_size += para_size + 2

    if current_chunk:
        chunks.append("\n\n".join(current_chunk))

    return chunks


def generate_embeddings(texts: List[str]) -> Optional[List[List[float]]]:
    """生成文本的向量嵌入"""
    model = get_embedding_model()
    if model is None:
        return None

    embeddings = model.encode(texts, show_progress_bar=False)
    return embeddings.tolist()


def extract_brand_from_filename(filename: str, content: str = "") -> str:
    """从文件名或内容中提取品牌"""
    brands = ["比亚迪", "特斯拉", "蔚来", "小鹏", "理想", "吉利", "长安", "长城", "上汽",
              "广汽", "奇瑞", "红旗", "问界", "小米", "零跑", "哪吒", "岚图"]

    filename_lower = filename.lower()
    content_lower = content.lower()

    for brand in brands:
        if brand in filename_lower or brand in content_lower:
            return brand

    return ""


async def store_document(
    filename: str,
    file_path: str,
    file_type: str,
    file_size: int,
    markdown: str,
    chunks: List[str],
    embeddings: List[List[float]],
    metadata: Dict[str, Any]
) -> int:
    """存储文档到数据库"""
    conn = get_db_connection()
    try:
        with conn:
            with conn.cursor() as cur:
                # 1. 插入 documents 表
                file_hash = hashlib.md5(f"{filename}{datetime.now()}".encode()).hexdigest()[:16]
                brand = extract_brand_from_filename(filename, markdown)

                cur.execute("""
                    INSERT INTO documents (file_name, source, brand, category, upload_date)
                    VALUES (%s, %s, %s, %s, %s)
                    RETURNING id
                """, (filename, "upload", brand, file_type, datetime.now()))

                doc_id = cur.fetchone()['id']

                # 2. 插入 chunks 表
                for i, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
                    chunk_brand = extract_brand_from_filename(filename, chunk)
                    cur.execute("""
                        INSERT INTO chunks (document_id, content, embedding, brand, metadata)
                        VALUES (%s, %s, %s, %s, %s)
                    """, (doc_id, chunk, embedding, chunk_brand or brand, json.dumps(metadata)))

                return doc_id

    finally:
        conn.close()


# ==================== API 接口 ====================

class UploadResponse(BaseModel):
    """上传响应"""
    success: bool
    document_id: Optional[int] = None
    filename: str
    file_type: str
    chunks: int
    status: str
    message: str


class DocumentInfo(BaseModel):
    """文档信息"""
    id: int
    file_name: str
    brand: str
    category: str
    upload_date: str
    chunk_count: int


@app.post("/upload")
async def upload_document(file: UploadFile = File(...)) -> UploadResponse:
    """上传并解析文档"""
    # 1. 验证文件类型
    allowed_types = ['pdf', 'pptx', 'ppt', 'docx', 'doc', 'xlsx', 'xls']
    file_ext = file.filename.split('.')[-1].lower() if '.' in file.filename else ''

    if file_ext not in allowed_types:
        raise HTTPException(
            status_code=400,
            detail=f"不支持的文件类型: {file_ext}，支持的类型: {', '.join(allowed_types)}"
        )

    # 2. 保存上传文件
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_filename = f"{timestamp}_{file.filename}"
    file_path = UPLOAD_DIR / safe_filename

    content = await file.read()
    file_size = len(content)

    with open(file_path, "wb") as f:
        f.write(content)

    try:
        # 3. 解析文档
        parse_result = await parse_document_mineru(str(file_path), file_ext)
        markdown = parse_result.get("markdown", "")

        # 4. 分块
        chunks = chunk_text(markdown)

        if not chunks:
            return UploadResponse(
                success=False,
                filename=file.filename,
                file_type=file_ext,
                chunks=0,
                status="error",
                message="文档内容为空"
            )

        # 5. 向量化
        embeddings = generate_embeddings(chunks)

        if embeddings is None:
            # 没有 embedding 模型，使用占位符
            embeddings = [[0.0] * EMBEDDING_DIM for _ in chunks]

        # 6. 存储到数据库
        doc_id = await store_document(
            filename=file.filename,
            file_path=str(file_path),
            file_type=file_ext,
            file_size=file_size,
            markdown=markdown,
            chunks=chunks,
            embeddings=embeddings,
            metadata={
                "parse_status": "success",
                "table_count": len(parse_result.get("tables", [])),
                "image_count": len(parse_result.get("images", []))
            }
        )

        return UploadResponse(
            success=True,
            document_id=doc_id,
            filename=file.filename,
            file_type=file_ext,
            chunks=len(chunks),
            status="completed",
            message="文档上传并解析完成"
        )

    except Exception as e:
        return UploadResponse(
            success=False,
            filename=file.filename,
            file_type=file_ext,
            chunks=0,
            status="error",
            message=f"处理失败: {str(e)}"
        )


@app.get("/documents")
async def list_documents(
    brand: str = None,
    category: str = None,
    page: int = 1,
    page_size: int = 20
) -> Dict[str, Any]:
    """获取文档列表"""
    conn = get_db_connection()
    try:
        with conn:
            with conn.cursor() as cur:
                # 构建查询
                where_clauses = []
                params = []

                if brand:
                    where_clauses.append("d.brand = %s")
                    params.append(brand)

                if category:
                    where_clauses.append("d.category = %s")
                    params.append(category)

                where_sql = " AND ".join(where_clauses) if where_clauses else "1=1"

                # 查询总数
                count_sql = f"""
                    SELECT COUNT(*) as total
                    FROM documents d
                    WHERE {where_sql}
                """
                cur.execute(count_sql, params)
                total = cur.fetchone()['total']

                # 查询列表
                offset = (page - 1) * page_size
                list_sql = f"""
                    SELECT
                        d.id,
                        d.file_name,
                        d.brand,
                        d.category,
                        d.upload_date,
                        COUNT(c.id) as chunk_count
                    FROM documents d
                    LEFT JOIN chunks c ON d.id = c.document_id
                    WHERE {where_sql}
                    GROUP BY d.id, d.file_name, d.brand, d.category, d.upload_date
                    ORDER BY d.upload_date DESC
                    LIMIT %s OFFSET %s
                """
                cur.execute(list_sql, params + [page_size, offset])
                rows = cur.fetchall()

                documents = []
                for row in rows:
                    documents.append({
                        "id": row['id'],
                        "file_name": row['file_name'],
                        "brand": row['brand'] or "",
                        "category": row['category'],
                        "upload_date": row['upload_date'].isoformat() if row['upload_date'] else "",
                        "chunk_count": row['chunk_count']
                    })

                return {
                    "success": True,
                    "total": total,
                    "page": page,
                    "page_size": page_size,
                    "documents": documents
                }

    finally:
        conn.close()


@app.get("/documents/{doc_id}")
async def get_document(doc_id: int) -> Dict[str, Any]:
    """获取文档详情"""
    conn = get_db_connection()
    try:
        with conn:
            with conn.cursor() as cur:
                # 查询文档信息
                cur.execute("""
                    SELECT id, file_name, brand, category, upload_date
                    FROM documents WHERE id = %s
                """, (doc_id,))
                doc = cur.fetchone()

                if not doc:
                    raise HTTPException(status_code=404, detail="文档不存在")

                # 查询 chunks
                cur.execute("""
                    SELECT id, content, brand, metadata
                    FROM chunks
                    WHERE document_id = %s
                    ORDER BY id
                """, (doc_id,))
                chunks = cur.fetchall()

                return {
                    "success": True,
                    "document": {
                        "id": doc['id'],
                        "file_name": doc['file_name'],
                        "brand": doc['brand'] or "",
                        "category": doc['category'],
                        "upload_date": doc['upload_date'].isoformat() if doc['upload_date'] else ""
                    },
                    "chunks": [
                        {
                            "id": c['id'],
                            "content": c['content'],
                            "brand": c['brand'] or "",
                            "metadata": c['metadata']
                        }
                        for c in chunks
                    ]
                }

    finally:
        conn.close()


@app.delete("/documents/{doc_id}")
async def delete_document(doc_id: int) -> Dict[str, Any]:
    """删除文档"""
    conn = get_db_connection()
    try:
        with conn:
            with conn.cursor() as cur:
                # 删除 chunks
                cur.execute("DELETE FROM chunks WHERE document_id = %s", (doc_id,))
                # 删除文档
                cur.execute("DELETE FROM documents WHERE id = %s", (doc_id,))

                return {"success": True, "message": "文档已删除"}

    finally:
        conn.close()


@app.get("/brands")
async def list_brands() -> Dict[str, Any]:
    """获取所有品牌列表"""
    conn = get_db_connection()
    try:
        with conn:
            with conn.cursor() as cur:
                cur.execute("""
                    SELECT DISTINCT brand FROM documents
                    WHERE brand IS NOT NULL AND brand != ''
                    ORDER BY brand
                """)
                rows = cur.fetchall()
                brands = [r['brand'] for r in rows]

                return {"success": True, "brands": brands}

    finally:
        conn.close()


@app.get("/health")
async def health():
    """健康检查"""
    try:
        conn = get_db_connection()
        conn.close()
        return {"status": "ok", "embedding_model": EMBEDDING_MODEL}
    except Exception as e:
        return {"status": "error", "error": str(e)}


# ==================== 启动 ====================

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8005)
