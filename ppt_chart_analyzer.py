# -*- coding: utf-8 -*-
"""
PPT图表数据分析工具 v1.0
功能：
1. PowerPoint COM自动化导出PPT页面为PNG图片
2. python-pptx提取图表数值数据
3. 提取PPT页面文字内容（表格、文本框）
4. 返回需要AI图像识别的图表列表（供image工具使用）
5. 生成结构化数据输出

依赖：pip install python-pptx Pillow pywin32
      需要安装Microsoft PowerPoint 或 WPS

作者：OpenClaw AI助手
日期：2026-04-01
"""

import sys
import os
import json
import zipfile
import xml.etree.ElementTree as ET
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple

# ============================
# 配置区
# ============================
WORKSPACE_DIR = r'E:\openclaw\workspace\ppt_analysis'
TEMP_IMAGE_DIR = WORKSPACE_DIR + r'\temp_images'
REPORT_OUTPUT_DIR = WORKSPACE_DIR + r'\reports'

# ============================
# PowerPoint COM 导出模块
# ============================

def export_ppt_to_images(ppt_path: str, output_dir: str = None, dpi: int = 1920) -> List[str]:
    """
    使用PowerPoint COM自动化将PPT导出为PNG图片
    
    Args:
        ppt_path: PPT文件路径
        output_dir: 输出目录，默认使用TEMP_IMAGE_DIR
        dpi: 导出分辨率，默认1920宽度
        
    Returns:
        导出的图片路径列表
    """
    import win32com.client
    
    if output_dir is None:
        output_dir = TEMP_IMAGE_DIR
    
    os.makedirs(output_dir, exist_ok=True)
    
    ppt_path = os.path.abspath(ppt_path)
    output_paths = []
    
    ppt = win32com.client.Dispatch('PowerPoint.Application')
    ppt.Visible = 1
    ppt.DisplayAlerts = 0  # ppAlertNone = 0
    
    try:
        prs = ppt.Presentations.Open(ppt_path, WithWindow=False)
        total_slides = prs.Slides.Count
        print(f'PPT总页数: {total_slides}')
        
        for i in range(1, total_slides + 1):
            slide = prs.Slides(i)
            out_path = os.path.join(output_dir, f'slide_{i:02d}.png')
            slide.Export(out_path, 'PNG')
            output_paths.append(out_path)
            print(f'  导出第{i}页: {out_path}')
        
        prs.Close()
        print(f'导出完成，共{len(output_paths)}页')
        
    except Exception as e:
        print(f'PowerPoint COM导出失败: {e}')
        raise
    finally:
        ppt.Quit()
    
    return output_paths


# ============================
# python-pptx 数值提取模块
# ============================

def extract_ppt_data(ppt_path: str) -> Dict[str, Any]:
    """
    使用python-pptx提取PPT中的所有数据
    
    Returns:
        {
            'slides': [
                {
                    'slide_index': 1,
                    'texts': [...],           # 所有文本内容
                    'tables': [...],           # 表格数据
                    'charts': [...],           # 图表数值数据
                    'images': [...],           # 图片列表
                    'has_chart': bool,
                    'image_needed': bool       # 是否需要图像识别补充
                },
                ...
            ],
            'warnings': [...]                 # 警告信息
        }
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    prs = Presentation(ppt_path)
    result = {
        'ppt_path': ppt_path,
        'ppt_name': os.path.basename(ppt_path),
        'total_slides': len(prs.slides),
        'slides': [],
        'warnings': []
    }
    
    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_data = {
            'slide_index': slide_idx,
            'texts': [],
            'tables': [],
            'charts': [],
            'images': [],
            'has_chart': False,
            'has_table': False,
            'image_needed': False
        }
        
        # 遍历所有shapes
        for shape in slide.shapes:
            # 文本内容
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                if len(text) < 500:  # 过滤超长文本
                    slide_data['texts'].append({
                        'text': text,
                        'shape_name': shape.name,
                        'position': {'left': shape.left, 'top': shape.top}
                    })
            
            # 表格
            if shape.has_table:
                slide_data['has_table'] = True
                table_data = extract_table_data(shape.table)
                slide_data['tables'].append(table_data)
            
            # 图表
            if shape.has_chart:
                slide_data['has_chart'] = True
                chart_data = extract_chart_data(shape.chart)
                slide_data['charts'].append(chart_data)
                
                # 检查是否需要图像识别（外链图表）
                if chart_data.get('is_external_link'):
                    slide_data['image_needed'] = True
                    result['warnings'].append(
                        f'第{slide_idx}页图表使用外链数据，需要图像识别补充标签'
                    )
        
        # 如果有图表但无法获取分类标签，标记需要图像识别
        if slide_data['has_chart'] and not any(
            c.get('categories') for c in slide_data['charts']
        ):
            slide_data['image_needed'] = True
        
        result['slides'].append(slide_data)
    
    return result


def extract_table_data(table) -> Dict[str, Any]:
    """提取表格数据"""
    rows = []
    for row in table.rows:
        row_data = [cell.text.strip() for cell in row.cells]
        rows.append(row_data)
    
    return {
        'rows': rows,
        'row_count': len(rows),
        'col_count': len(rows[0]) if rows else 0
    }


def extract_chart_data(chart) -> Dict[str, Any]:
    """提取图表数据"""
    chart_info = {
        'chart_type': str(chart.chart_type),
        'series': [],
        'categories': [],  # 可能为空（外链图表）
        'values': [],
        'is_external_link': False
    }
    
    ns = {
        'c': 'http://schemas.openxmlformats.org/drawingml/2006/chart',
    }
    
    try:
        with zipfile.ZipFile(extract_pptx_path()) as z:
            # 找到对应的chart XML
            chart_files = [n for n in z.namelist() if 'chart' in n and n.endswith('.xml')]
            # 简化处理：通过series数量匹配
    except:
        pass
    
    # 提取series数据
    for series in chart.series:
        series_data = {
            'name': series.name if hasattr(series, 'name') else f'Series{len(chart_info["series"])+1}',
            'values': list(series.values) if hasattr(series, 'values') else []
        }
        chart_info['series'].append(series_data)
        chart_info['values'].extend(series_data['values'])
    
    # 检查是否需要图像识别（categories为空但有多个series）
    if not chart_info['categories'] and len(chart_info['series']) > 0:
        chart_info['is_external_link'] = True
    
    return chart_info


def extract_pptx_path():
    """获取当前PPT路径（临时方法）"""
    return TEMP_IMAGE_DIR  # 临时占位


# ============================
# 图像识别任务生成模块
# ============================

def generate_image_analysis_prompt(slide_data: Dict, chart_index: int = 0) -> str:
    """
    生成图像分析的提示词，供image工具使用
    
    Args:
        slide_data: 单页幻灯片数据
        chart_index: 要分析的图表索引
        
    Returns:
        分析提示词
    """
    chart_data = slide_data['charts'][chart_index] if slide_data['charts'] else {}
    
    prompt = f"""请仔细分析这张PPT幻灯片截图中的图表内容：

1. **图表类型**: {chart_data.get('chart_type', '未知')}
2. **已有数值数据**: {chart_data.get('values', [])}
3. **请识别**:
   - X轴分类标签（如：轿车、SUV、MPV等）
   - Y轴标签和单位
   - 图例中各颜色对应的分类名称
   - 各数据系列的名称和颜色对应关系
   - 图表中显示的具体数值标注

请完整提取所有可识别的文字和数据信息。"""
    
    return prompt


def build_analysis_report(data: Dict, web_insights: str = "") -> str:
    """
    生成结构化分析报告
    
    Args:
        data: extract_ppt_data() 的输出
        web_insights: 网络搜索获取的市场洞察
        
    Returns:
        Markdown格式的分析报告
    """
    report = []
    report.append(f"# {data['ppt_name']} 数据分析报告")
    report.append(f"\n**生成时间**: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    report.append(f"**分析页数**: {data['total_slides']} 页")
    
    # 警告信息
    if data['warnings']:
        report.append("\n## ⚠️ 数据质量警告")
        for warning in data['warnings']:
            report.append(f"- {warning}")
    
    # 每页数据分析
    for slide in data['slides']:
        slide_num = slide['slide_index']
        report.append(f"\n---\n\n## 第 {slide_num} 页")
        
        # 文本内容
        if slide['texts']:
            report.append("\n### 文本内容")
            for t in slide['texts'][:10]:  # 限制数量
                if len(t['text']) > 2:
                    report.append(f"- {t['text']}")
        
        # 表格数据
        if slide['tables']:
            report.append("\n### 表格数据")
            for tbl in slide['tables']:
                report.append(f"| {' | '.join(tbl['rows'][0])} |")
                report.append(f"| {' | '.join(['---'] * tbl['col_count'])} |")
                for row in tbl['rows'][1:5]:
                    report.append(f"| {' | '.join(row)} |")
        
        # 图表数据
        if slide['charts']:
            report.append("\n### 图表数据")
            for i, chart in enumerate(slide['charts']):
                report.append(f"\n**图表 {i+1}** (类型: {chart['chart_type']})")
                
                if chart['categories']:
                    report.append(f"- 分类: {', '.join(str(c) for c in chart['categories'])}")
                
                for j, series in enumerate(chart['series']):
                    report.append(f"- 系列{j+1} ({series['name']}): {series['values']}")
                
                if chart['is_external_link']:
                    report.append("\n⚠️ **注意**: 此图表使用外链数据，分类标签需要通过图像识别补充")
        
        # 需要图像识别
        if slide['image_needed']:
            report.append("\n### 🤖 图像识别任务")
            prompt = generate_image_analysis_prompt(slide)
            report.append("\n请使用image工具分析对应截图，提示词：")
            report.append(f"\n```\n{prompt}\n```")
    
    # 市场洞察
    if web_insights:
        report.append("\n---\n\n## 🌐 市场背景与深层洞察")
        report.append(f"\n{web_insights}")
    
    return '\n'.join(report)


# ============================
# 主工作流
# ============================

def analyze_ppt_workflow(ppt_path: str, output_dir: str = None) -> Dict[str, Any]:
    """
    完整工作流：PPT文件 → 图片导出 → 数据提取 → 报告生成
    
    Args:
        ppt_path: PPT文件路径
        output_dir: 输出目录
        
    Returns:
        {
            'ppt_data': extract_ppt_data()的结果,
            'image_paths': 导出的图片路径列表,
            'report': 生成的报告内容,
            'image_analysis_tasks': 需要图像识别的任务列表
        }
    """
    print(f'\n{"="*60}')
    print(f'PPT分析工作流启动: {ppt_path}')
    print(f'{"="*60}\n')
    
    # Step 1: 导出PPT页面为图片
    print('[Step 1/4] 导出PPT页面为PNG图片...')
    image_paths = export_ppt_to_images(ppt_path, output_dir)
    print(f'导出完成: {len(image_paths)} 页\n')
    
    # Step 2: 提取PPT数值数据
    print('[Step 2/4] 提取PPT数值数据...')
    ppt_data = extract_ppt_data(ppt_path)
    slide_with_charts = sum(1 for s in ppt_data['slides'] if s['has_chart'])
    print(f'发现 {slide_with_charts} 页包含图表\n')
    
    # Step 3: 生成图像分析任务列表
    print('[Step 3/4] 生成图像分析任务...')
    image_tasks = []
    for slide in ppt_data['slides']:
        if slide['image_needed']:
            image_path = image_paths[slide['slide_index'] - 1]
            prompt = generate_image_analysis_prompt(slide)
            image_tasks.append({
                'slide_index': slide['slide_index'],
                'image_path': image_path,
                'prompt': prompt
            })
    print(f'需要图像识别: {len(image_tasks)} 个任务\n')
    
    # Step 4: 生成初步报告
    print('[Step 4/4] 生成结构化报告...')
    report = build_analysis_report(ppt_data)
    print('报告生成完成\n')
    
    # 保存报告
    os.makedirs(REPORT_OUTPUT_DIR, exist_ok=True)
    report_path = os.path.join(REPORT_OUTPUT_DIR, 
                               f'report_{datetime.now().strftime("%Y%m%d_%H%M")}.md')
    with open(report_path, 'w', encoding='utf-8') as f:
        f.write(report)
    print(f'报告已保存: {report_path}')
    
    return {
        'ppt_data': ppt_data,
        'image_paths': image_paths,
        'report': report,
        'report_path': report_path,
        'image_analysis_tasks': image_tasks
    }


# ============================
# 使用示例
# ============================

if __name__ == '__main__':
    # 示例用法
    print("""
PPT图表数据分析工具 v1.0
========================

使用方法:
1. 导入: from ppt_chart_analyzer import analyze_ppt_workflow
2. 调用: result = analyze_ppt_workflow('your_file.pptx')
3. 查看图像任务: result['image_analysis_tasks']
4. 使用image工具分析每张图片
5. 合并结果生成最终报告

依赖安装:
    pip install python-pptx Pillow pywin32
""")
