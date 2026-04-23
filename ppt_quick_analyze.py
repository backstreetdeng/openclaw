# -*- coding: utf-8 -*-
"""
PPT快速分析工具 v1.0 - 简洁版
一键执行完整工作流：PPT → 导出图片 → 提取数据 → 图像识别 → 生成报告

依赖：pip install python-pptx Pillow pywin32
      需要安装 Microsoft PowerPoint 或 WPS
"""

import sys
import os
import json
import zipfile
import re
from datetime import datetime
from typing import Dict, List, Any

# ============================
# 核心配置
# ============================
WORKSPACE = r'E:\openclaw\workspace\ppt_analysis'
TEMP_IMAGES = WORKSPACE + r'\temp_images'
REPORTS = WORKSPACE + r'\reports'

# ============================
# Step 1: PowerPoint导出PNG
# ============================

def export_slides(ppt_path: str) -> List[str]:
    """将PPT每一页导出为PNG图片"""
    import win32com.client
    
    os.makedirs(TEMP_IMAGES, exist_ok=True)
    ppt_path = os.path.abspath(ppt_path)
    paths = []
    
    ppt = win32com.client.Dispatch('PowerPoint.Application')
    ppt.Visible = 1
    ppt.DisplayAlerts = 0
    
    try:
        prs = ppt.Presentations.Open(ppt_path, WithWindow=False)
        print(f'PPT总页数: {prs.Slides.Count}')
        
        for i in range(1, prs.Slides.Count + 1):
            out = os.path.join(TEMP_IMAGES, f's_{i:02d}.png')
            prs.Slides(i).Export(out, 'PNG')
            paths.append(out)
            print(f'  导出第{i}页')
        
        prs.Close()
    finally:
        ppt.Quit()
    
    return paths


# ============================
# Step 2: python-pptx提取数据
# ============================

def extract_data(ppt_path: str) -> Dict[str, Any]:
    """提取PPT中所有图表和表格的数值数据"""
    from pptx import Presentation
    
    prs = Presentation(ppt_path)
    result = {'slides': [], 'warnings': []}
    
    for idx, slide in enumerate(prs.slides, 1):
        slide_data = {
            'index': idx,
            'texts': [],
            'tables': [],
            'charts': [],
            'need_image': False
        }
        
        for shape in slide.shapes:
            # 文本
            if hasattr(shape, 'text') and shape.text.strip():
                txt = shape.text.strip()
                if 2 < len(txt) < 200:
                    slide_data['texts'].append(txt)
            
            # 表格
            if shape.has_table:
                rows = []
                for row in shape.table.rows:
                    rows.append([c.text.strip() for c in row.cells])
                slide_data['tables'].append(rows)
            
            # 图表
            if shape.has_chart:
                chart_vals = []
                for series in shape.chart.series:
                    vals = list(series.values) if hasattr(series, 'values') else []
                    chart_vals.append(vals)
                slide_data['charts'].append({
                    'type': str(shape.chart.chart_type),
                    'series': chart_vals
                })
        
        # 检查是否需要图像识别：图表有无分类标签
        # 如果图表series数量>0但没有文本标签（标题/表头）包含车型名称，则需要图像识别
        if slide_data['charts']:
            # 常见车型关键词
            car_keywords = ['轿车', 'SUV', 'MPV', '交叉型', '乘用车', '新能源', '燃油', '纯电', '插电', '混动']
            has_car_labels = any(
                any(kw in t for kw in car_keywords) 
                for t in slide_data['texts']
            )
            if not has_car_labels:
                slide_data['need_image'] = True
        
        result['slides'].append(slide_data)
    
    return result


# ============================
# Step 3: 图像识别（调用image工具）
# ============================

def analyze_image(image_path: str, chart_info: str = "") -> str:
    """
    使用image工具分析图片
    注意：此函数需要AI代理在分析时调用image工具
    这里只生成提示词
    """
    prompt = f"""请仔细分析这张PPT幻灯片截图：

1. 识别所有文字内容（标题、标签、标注）
2. 识别图表的X轴分类名称（如：轿车、SUV、MPV等）
3. 识别图例中各颜色对应的分类
4. 提取表格中的所有数据
5. 如有图表，提取各柱状/曲线的具体数值

请以结构化方式输出所有信息。"""
    return prompt


# ============================
# Step 4: 生成报告
# ============================

def generate_report(ppt_name: str, data: Dict, image_results: Dict = None) -> str:
    """生成Markdown分析报告"""
    lines = []
    lines.append(f"# {ppt_name} 数据分析报告")
    lines.append(f"\n**生成时间**: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    lines.append(f"**总页数**: {len(data['slides'])}")
    
    for slide in data['slides']:
        lines.append(f"\n---")
        lines.append(f"\n## 第 {slide['index']} 页")
        
        if slide['texts']:
            lines.append("\n**文本内容**")
            for t in slide['texts'][:15]:
                lines.append(f"- {t}")
        
        if slide['tables']:
            lines.append("\n**表格数据**")
            for tbl in slide['tables']:
                for row in tbl[:8]:
                    lines.append(f"| {' | '.join(row)} |")
        
        if slide['charts']:
            lines.append("\n**图表数据**")
            for i, ch in enumerate(slide['charts']):
                lines.append(f"- 类型: {ch['type']}")
                for j, vals in enumerate(ch['series']):
                    lines.append(f"  - 系列{j+1}: {vals[:10]}{'...' if len(vals)>10 else ''}")
        
        if slide['need_image']:
            lines.append("\n⚠️ **需要图像识别**: 请使用image工具分析 `temp_images/s_{:02d}.png`")
    
    if image_results:
        lines.append("\n---\n\n## 图像识别补充结果")
        for page, result in image_results.items():
            lines.append(f"\n### 第{page}页")
            lines.append(result)
    
    return '\n'.join(lines)


# ============================
# 完整工作流
# ============================

def run(ppt_path: str) -> Dict[str, Any]:
    """
    执行完整分析工作流
    
    Returns:
        {
            'image_paths': List[str],    # 导出的图片路径
            'data': Dict,                # 提取的数值数据
            'report': str,               # 生成的报告
            'image_tasks': List[Dict]    # 需要image工具分析的任务
        }
    """
    print(f'\n{"="*50}')
    print(f'开始分析: {ppt_path}')
    print(f'{"="*50}\n')
    
    # Step 1: 导出图片
    print('[1/4] 导出PPT页面为PNG...')
    image_paths = export_slides(ppt_path)
    print(f'完成: {len(image_paths)} 页\n')
    
    # Step 2: 提取数据
    print('[2/4] 提取PPT数值数据...')
    data = extract_data(ppt_path)
    charts_count = sum(len(s['charts']) for s in data['slides'])
    print(f'完成: 发现 {charts_count} 个图表\n')
    
    # Step 3: 生成图像分析任务
    print('[3/4] 生成图像分析任务...')
    image_tasks = []
    for slide in data['slides']:
        if slide['need_image']:
            img_path = image_paths[slide['index'] - 1]
            image_tasks.append({
                'page': slide['index'],
                'image_path': img_path,
                'prompt': analyze_image(img_path)
            })
    print(f'完成: {len(image_tasks)} 个任务需要图像识别\n')
    
    # Step 4: 生成报告
    print('[4/4] 生成分析报告...')
    report = generate_report(os.path.basename(ppt_path), data)
    
    os.makedirs(REPORTS, exist_ok=True)
    report_path = os.path.join(REPORTS, f'report_{datetime.now().strftime("%Y%m%d_%H%M")}.md')
    with open(report_path, 'w', encoding='utf-8') as f:
        f.write(report)
    print(f'报告已保存: {report_path}\n')
    
    print(f'{"="*50}')
    print('分析完成!')
    print(f'{"="*50}')
    print(f'\n图片目录: {TEMP_IMAGES}')
    print(f'报告路径: {report_path}')
    print(f'\n下一步: 使用image工具分析以下页面:')
    for task in image_tasks:
        print(f"  - 第{task['page']}页: {task['image_path']}")
    
    return {
        'image_paths': image_paths,
        'data': data,
        'report': report,
        'report_path': report_path,
        'image_tasks': image_tasks
    }


# ============================
# 使用说明
# ============================

if __name__ == '__main__':
    print("""
╔══════════════════════════════════════════════════════════════╗
║         PPT图表数据分析工具 v1.0 - 使用说明                  ║
╠══════════════════════════════════════════════════════════════╣
║                                                              ║
║  使用方法:                                                   ║
║                                                              ║
║  from ppt_quick_analyze import run                           ║
║  result = run(r'E:\\your\\file.pptx')                       ║
║                                                              ║
║  返回结果:                                                   ║
║  - result['image_paths']   : 导出的图片路径列表              ║
║  - result['data']          : 提取的数值数据                  ║
║  - result['report']        : 初步分析报告                    ║
║  - result['image_tasks']   : 需要image工具分析的任务         ║
║  - result['report_path']   : 报告保存路径                    ║
║                                                              ║
║  完整工作流:                                                  ║
║  1. run() 导出图片 + 提取数据 + 生成初步报告                 ║
║  2. 对每个 image_task 用 image工具 分析对应图片               ║
║  3. 将图像识别结果补充到报告中                               ║
║  4. 结合网络搜索的市场洞察，生成最终报告                      ║
║                                                              ║
║  依赖: pip install python-pptx Pillow pywin32               ║
║                                                              ║
╚══════════════════════════════════════════════════════════════╝
""")
